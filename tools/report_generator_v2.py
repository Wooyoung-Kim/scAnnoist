"""
Enhanced Report Generator using HTML2PPTX workflow
Generates professional presentations from scRNA-seq analysis results
"""

import json
import subprocess
import base64
from pathlib import Path
from typing import Dict, Any, List, Optional
import pandas as pd


class PresentationGeneratorV2:
    """
    Generate presentations using HTML2PPTX workflow.
    Much faster and more reliable than AI-based generation.
    """

    def __init__(self, output_dir: str):
        """
        Initialize generator.

        Args:
            output_dir: Directory containing analysis results
        """
        self.output_dir = Path(output_dir)
        self.slides_dir = self.output_dir / "slides_html"
        self.slides_dir.mkdir(exist_ok=True, parents=True)

        # Find figures directory
        self.figures_dir = self.output_dir.parent / "figures"

        # Find html2pptx CLI script
        self.html2pptx_script = Path("/mnt2/kwy/scrna_agent/.claude/skills/pptx/scripts/html2pptx_cli.js")

        if not self.html2pptx_script.exists():
            raise FileNotFoundError(f"html2pptx_cli.js not found at {self.html2pptx_script}")

    def load_analysis_results(self) -> Dict[str, Any]:
        """Load analysis results from output directory."""
        results = {}

        # Load predictions
        pred_file = self.output_dir / "celltypist_predictions.csv"
        if pred_file.exists():
            results['predictions'] = pd.read_csv(pred_file)

        # Load tissue usage
        tissue_file = self.output_dir / "tissue_model_usage.csv"
        if tissue_file.exists():
            results['tissue_usage'] = pd.read_csv(tissue_file)

        # Load metadata
        meta_file = self.output_dir / "run_metadata.json"
        if meta_file.exists():
            with open(meta_file) as f:
                results['metadata'] = json.load(f)

        return results

    def generate_html_slide(self, slide_num: int, title: str, content_html: str) -> Path:
        """
        Generate a single HTML slide.

        Args:
            slide_num: Slide number
            title: Slide title
            content_html: HTML content for the slide body

        Returns:
            Path to HTML file
        """
        html_file = self.slides_dir / f"slide_{slide_num:02d}.html"

        # Color palette - Scientific Blue & Teal
        primary_color = "#1C2833"  # Deep navy
        secondary_color = "#5EA8A7"  # Teal
        accent_color = "#E74C3C"  # Coral red for emphasis
        bg_color = "#F4F6F6"  # Off-white

        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        body {{
            width: 720pt;
            height: 405pt;
            display: flex;
            flex-direction: column;
            font-family: Arial, sans-serif;
            background: {bg_color};
            overflow: hidden;
        }}
        .header {{
            background: {primary_color};
            padding: 15pt 40pt;
            border-left: 8pt solid {secondary_color};
        }}
        .header h1 {{
            color: #FFFFFF;
            font-size: 24pt;
            font-weight: bold;
        }}
        .content {{
            flex: 1;
            padding: 20pt 40pt 36pt 40pt;
            display: flex;
            flex-direction: column;
            gap: 12pt;
        }}
        h2 {{
            color: {primary_color};
            font-size: 20pt;
            font-weight: bold;
            margin-bottom: 8pt;
        }}
        h3 {{
            color: {secondary_color};
            font-size: 16pt;
            font-weight: bold;
            margin-bottom: 6pt;
        }}
        p {{
            color: #2C2C2C;
            font-size: 12pt;
            line-height: 1.4;
        }}
        ul {{
            color: #2C2C2C;
            font-size: 12pt;
            line-height: 1.4;
            margin-left: 20pt;
        }}
        .stat-box {{
            background: #FFFFFF;
            border-left: 4pt solid {secondary_color};
            padding: 12pt 16pt;
            border-radius: 4pt;
            box-shadow: 2px 2px 8px rgba(0, 0, 0, 0.1);
        }}
        .stat-number {{
            color: {accent_color};
            font-size: 24pt;
            font-weight: bold;
        }}
        .stat-label {{
            color: #666666;
            font-size: 10pt;
        }}
        .footer {{
            padding: 10pt 40pt;
            text-align: right;
        }}
        .footer p {{
            color: #999999;
            font-size: 10pt;
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{title}</h1>
    </div>
    <div class="content">
        {content_html}
    </div>
    <div class="footer">
        <p>Generated with scRNA-agent | Slide {slide_num}</p>
    </div>
</body>
</html>"""

        with open(html_file, 'w', encoding='utf-8') as f:
            f.write(html_content)

        return html_file

    def create_title_slide(self, title: str) -> Path:
        """Create title slide."""
        content = f"""
        <div style="display: flex; flex-direction: column; justify-content: center; align-items: center; height: 100%; gap: 20pt;">
            <h2 style="font-size: 36pt; text-align: center; color: #1C2833;">{title}</h2>
            <p style="font-size: 18pt; text-align: center; color: #5EA8A7;">Tissue-Aware Cell Type Annotation</p>
            <p style="font-size: 14pt; text-align: center; color: #666666; margin-top: 30pt;">Powered by CellTypist &amp; scRNA-agent</p>
        </div>
        """
        return self.generate_html_slide(1, "", content)

    def create_overview_slide(self, results: Dict) -> Path:
        """Create analysis overview slide."""
        meta = results.get('metadata', {})
        total_cells = meta.get('total_cells', 0)
        total_tissues = meta.get('total_tissues', 0)
        cell_types = meta.get('cell_types_identified', 0)

        content = f"""
        <h2>Analysis Overview</h2>
        <div style="display: flex; gap: 15pt; margin-top: 15pt;">
            <div class="stat-box" style="flex: 1;">
                <p class="stat-number">{total_cells:,}</p>
                <p class="stat-label">Cells Analyzed</p>
            </div>
            <div class="stat-box" style="flex: 1;">
                <p class="stat-number">{total_tissues}</p>
                <p class="stat-label">Tissue Types</p>
            </div>
            <div class="stat-box" style="flex: 1;">
                <p class="stat-number">{cell_types}</p>
                <p class="stat-label">Cell Types Identified</p>
            </div>
        </div>
        <div style="margin-top: 20pt;">
            <h3>Annotation Method</h3>
            <p>CellTypist with tissue-specific model selection. Majority voting enabled for enhanced accuracy.</p>
        </div>
        """
        return self.generate_html_slide(2, "Analysis Overview", content)

    def create_cell_types_slide(self, results: Dict) -> Path:
        """Create cell types distribution slide."""
        meta = results.get('metadata', {})
        top_types = meta.get('top_cell_types', {})

        # Get top 8 (to fit in slide)
        items = list(top_types.items())[:8]
        total = sum(top_types.values())

        rows_html = ""
        for cell_type, count in items:
            percentage = (count / total) * 100
            bar_width = percentage * 2.5  # Scale for visualization
            rows_html += f"""
            <div style="margin-bottom: 8pt;">
                <p style="font-weight: bold; color: #1C2833; margin-bottom: 3pt; font-size: 11pt;">{cell_type}</p>
                <div style="display: flex; align-items: center; gap: 8pt;">
                    <div style="width: 250pt; height: 14pt; background: #E0E0E0; border-radius: 7pt;">
                        <div style="width: {bar_width}pt; height: 100%; background: #5EA8A7; border-radius: 7pt;"></div>
                    </div>
                    <p style="font-size: 10pt; color: #666666;">{count:,} ({percentage:.1f}%)</p>
                </div>
            </div>
            """

        content = f"""
        <h2>Top Cell Types Identified</h2>
        <div style="margin-top: 15pt;">
            {rows_html}
        </div>
        """
        return self.generate_html_slide(3, "Cell Type Distribution", content)

    def create_tissue_info_slide(self, results: Dict) -> Path:
        """Create tissue and model information slide."""
        meta = results.get('metadata', {})
        tissue_mapping = meta.get('tissue_model_mapping', {})

        rows_html = ""
        for tissue, info in tissue_mapping.items():
            model = info.get('model', 'Unknown')
            rule = info.get('rule', 'Unknown')
            rows_html += f"""
            <div class="stat-box" style="margin-bottom: 12pt;">
                <h3 style="color: #5EA8A7; margin-bottom: 6pt; font-size: 14pt;">{tissue}</h3>
                <p style="font-size: 11pt;"><b>Model:</b> {model}</p>
                <p style="font-size: 10pt; color: #666666; margin-top: 3pt;">Selection rule: {rule}</p>
            </div>
            """

        content = f"""
        <h2>Model Selection per Tissue</h2>
        <div style="margin-top: 15pt;">
            {rows_html}
        </div>
        <div style="margin-top: 15pt;">
            <p style="font-size: 11pt; color: #666666;">Models automatically selected based on tissue type and sample characteristics</p>
        </div>
        """
        return self.generate_html_slide(4, "Models Used", content)

    def create_umap_cluster_slide(self) -> Optional[Path]:
        """Create UMAP cluster visualization slide."""
        img_path = self.figures_dir / "umap_clusters.png"
        if not img_path.exists():
            return None

        # Encode image to base64
        with open(img_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode()

        content = f"""
        <div style="display: flex; flex-direction: column; align-items: center; gap: 15pt;">
            <img src="data:image/png;base64,{img_data}" style="width: 550pt; height: auto;" />
            <div style="margin-top: 10pt; padding: 0 20pt;">
                <h3>Interpretation</h3>
                <p>UMAP projection reveals distinct cell populations based on gene expression patterns. Each color represents a Leiden cluster, showing the natural grouping of cells with similar transcriptional profiles.</p>
            </div>
        </div>
        """
        return self.generate_html_slide(3, "Clustering Results", content)

    def create_umap_celltype_slide(self) -> Optional[Path]:
        """Create UMAP cell type visualization slide."""
        img_path = self.figures_dir / "umap_celltypes.png"
        if not img_path.exists():
            return None

        # Encode image to base64
        with open(img_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode()

        content = f"""
        <div style="display: flex; flex-direction: column; align-items: center; gap: 15pt;">
            <img src="data:image/png;base64,{img_data}" style="width: 550pt; height: auto;" />
            <div style="margin-top: 10pt; padding: 0 20pt;">
                <h3>Interpretation</h3>
                <p>Cell type annotations from CellTypist mapped onto UMAP coordinates. Spatial segregation of cell types validates the quality of both clustering and annotation. Each cell is colored by its predicted cell type identity.</p>
            </div>
        </div>
        """
        return self.generate_html_slide(4, "CellTypist Annotations", content)

    def create_marker_annotation_slide(self) -> Optional[Path]:
        """Create marker-based annotation visualization slide."""
        img_path = self.figures_dir / "umap_marker.png"
        if not img_path.exists():
            return None

        # Encode image to base64
        with open(img_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode()

        content = f"""
        <div style="display: flex; flex-direction: column; align-items: center; gap: 15pt;">
            <img src="data:image/png;base64,{img_data}" style="width: 550pt; height: auto;" />
            <div style="margin-top: 10pt; padding: 0 20pt;">
                <h3>Interpretation</h3>
                <p>Marker-based cell type annotations using known marker genes for each cell type. This method provides independent validation of CellTypist results and identifies cell types based on canonical marker expression patterns.</p>
            </div>
        </div>
        """
        return self.generate_html_slide(5, "Marker-Based Annotations", content)

    def create_annotation_comparison_slide(self) -> Optional[Path]:
        """Create annotation method comparison slide."""
        img_path = self.figures_dir / "annotation_comparison.png"
        if not img_path.exists():
            return None

        # Encode image to base64
        with open(img_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode()

        content = f"""
        <div style="display: flex; flex-direction: column; align-items: center; gap: 12pt;">
            <img src="data:image/png;base64,{img_data}" style="width: 620pt; height: auto;" />
            <div style="margin-top: 8pt; padding: 0 20pt;">
                <h3>Interpretation</h3>
                <p>Side-by-side comparison of CellTypist (left) and marker-based (right) annotations. Agreement between methods validates annotation quality, while discrepancies may indicate rare cell types or transitional states requiring further investigation.</p>
            </div>
        </div>
        """
        return self.generate_html_slide(6, "Method Comparison", content)

    def create_qc_slide(self) -> Optional[Path]:
        """Create QC metrics visualization slide."""
        img_path = self.figures_dir / "qc_metrics.png"
        if not img_path.exists():
            return None

        # Encode image to base64
        with open(img_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode()

        content = f"""
        <div style="display: flex; flex-direction: column; align-items: center; gap: 15pt;">
            <img src="data:image/png;base64,{img_data}" style="width: 600pt; height: auto;" />
            <div style="margin-top: 5pt; padding: 0 20pt;">
                <h3>Interpretation</h3>
                <p>Quality control metrics show number of genes detected, total counts, and mitochondrial content per cell. Filtering was applied to remove low-quality cells and ensure high-confidence downstream analysis.</p>
            </div>
        </div>
        """
        return self.generate_html_slide(2, "Quality Control", content)

    def create_summary_slide(self, results: Dict) -> Path:
        """Create summary slide."""
        meta = results.get('metadata', {})
        total_cells = meta.get('total_cells', 0)
        cell_types = meta.get('cell_types_identified', 0)

        content = f"""
        <h2>Summary</h2>
        <div style="margin-top: 15pt;">
            <h3>Key Findings</h3>
            <ul>
                <li>Successfully annotated {total_cells:,} cells</li>
                <li>Identified {cell_types} distinct cell types</li>
                <li>Tissue-specific models ensured accurate classification</li>
            </ul>
        </div>
        <div style="margin-top: 20pt;">
            <h3>Next Steps</h3>
            <ul>
                <li>Perform differential expression analysis</li>
                <li>Investigate cell-cell interactions</li>
                <li>Validate with marker gene expression</li>
            </ul>
        </div>
        """
        return self.generate_html_slide(99, "Summary & Next Steps", content)

    def convert_html_to_pptx(self, html_files: List[Path], output_pptx: Path) -> None:
        """
        Convert HTML slides to PPTX using html2pptx.js.

        Args:
            html_files: List of HTML slide files
            output_pptx: Output PPTX file path
        """
        print(f"\nConverting {len(html_files)} HTML slides to PowerPoint...")

        for i, html_file in enumerate(html_files, 1):
            print(f"  Converting slide {i}/{len(html_files)}...")

            # Output PPTX for this slide
            slide_pptx = self.slides_dir / f"slide_{i:02d}.pptx"

            # Run html2pptx.js
            cmd = [
                "node",
                str(self.html2pptx_script),
                str(html_file),
                str(slide_pptx)
            ]

            try:
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=30
                )

                if result.returncode != 0:
                    print(f"    Error: {result.stderr}")
                    continue

                # Check if file was created
                if slide_pptx.exists():
                    print(f"    ✓ Created {slide_pptx.name}")
                else:
                    print(f"    Warning: {slide_pptx.name} not found")

            except Exception as e:
                print(f"    Error: {e}")

        # Combine all PPTX files
        self.combine_pptx_files(output_pptx)

    def combine_pptx_files(self, output_pptx: Path) -> None:
        """
        Combine individual PPTX files into one presentation.

        Args:
            output_pptx: Output PPTX file path
        """
        print(f"\nCombining slides into final presentation...")

        # Get all slide PPTX files
        slide_files = sorted(self.slides_dir.glob("slide_*.pptx"))

        if not slide_files:
            raise ValueError("No PPTX slides found to combine")

        # Use first file as base, we'll manually merge using python-pptx if available
        # For now, just use the first slide as the output (single slide presentation)
        # In production, would use proper PPTX merging library

        # Simple approach: just copy the first one for now
        # TODO: Implement proper multi-slide PPTX merging
        import shutil
        if len(slide_files) == 1:
            shutil.copy(slide_files[0], output_pptx)
        else:
            # For multiple files, we'll create a simple combined version
            # This requires python-pptx library
            try:
                from pptx import Presentation
                from pptx.util import Inches

                # Start with blank presentation
                prs = Presentation()
                prs.slide_width = Inches(10)  # 720pt = 10 inches
                prs.slide_height = Inches(5.625)  # 405pt = 5.625 inches

                # Load each slide file and copy slides
                for slide_file in slide_files:
                    src_prs = Presentation(str(slide_file))
                    for slide in src_prs.slides:
                        # Add blank slide
                        blank_layout = prs.slide_layouts[6]  # Blank layout
                        new_slide = prs.slides.add_slide(blank_layout)

                        # Copy all shapes from source slide
                        for shape in slide.shapes:
                            # This is a simplified copy - may not preserve all formatting
                            try:
                                el = shape.element
                                new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                            except:
                                pass  # Skip if copy fails

                prs.save(str(output_pptx))
                print(f"✓ Saved combined presentation: {output_pptx}")

            except ImportError:
                # If python-pptx not available, just use first slide
                print("Warning: python-pptx not available, using first slide only")
                shutil.copy(slide_files[0], output_pptx)

    def generate_presentation(self, title: str) -> Path:
        """
        Generate complete presentation.

        Args:
            title: Presentation title

        Returns:
            Path to output PDF/PPTX file
        """
        print("=" * 60)
        print("Generating Presentation (HTML2PPTX)")
        print("=" * 60)

        # Load results
        print("Loading analysis results...")
        results = self.load_analysis_results()

        if not results:
            raise ValueError("No analysis results found")

        # Generate HTML slides
        print("Creating HTML slides...")
        html_files = []

        # Title slide
        html_files.append(self.create_title_slide(title))

        # QC slide (if available)
        qc_slide = self.create_qc_slide()
        if qc_slide:
            html_files.append(qc_slide)
            print(f"  ✓ Added QC metrics slide")

        # Clustering UMAP slide (if available)
        cluster_slide = self.create_umap_cluster_slide()
        if cluster_slide:
            html_files.append(cluster_slide)
            print(f"  ✓ Added clustering UMAP slide")

        # Cell type UMAP slide (if available)
        celltype_slide = self.create_umap_celltype_slide()
        if celltype_slide:
            html_files.append(celltype_slide)
            print(f"  ✓ Added CellTypist UMAP slide")

        # Marker-based annotation slide (if available)
        marker_slide = self.create_marker_annotation_slide()
        if marker_slide:
            html_files.append(marker_slide)
            print(f"  ✓ Added marker-based annotation slide")

        # Annotation comparison slide (if available)
        comparison_slide = self.create_annotation_comparison_slide()
        if comparison_slide:
            html_files.append(comparison_slide)
            print(f"  ✓ Added annotation comparison slide")

        # Overview and results slides
        html_files.append(self.create_overview_slide(results))
        html_files.append(self.create_cell_types_slide(results))
        html_files.append(self.create_tissue_info_slide(results))
        html_files.append(self.create_summary_slide(results))

        print(f"✓ Created {len(html_files)} HTML slides total")

        # Convert to PPTX
        output_pptx = self.output_dir / "presentation.pptx"
        self.convert_html_to_pptx(html_files, output_pptx)

        return output_pptx


def generate_report_v2(output_dir: str, title: str = "scRNA-seq Analysis Results") -> Path:
    """
    Generate presentation report using HTML2PPTX workflow.

    Args:
        output_dir: Directory containing analysis results
        title: Presentation title

    Returns:
        Path to generated PPTX file
    """
    generator = PresentationGeneratorV2(output_dir)
    return generator.generate_presentation(title)
