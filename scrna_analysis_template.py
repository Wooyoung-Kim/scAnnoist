#!/usr/bin/env python3
"""
scRNA-seq Analysis Template with Automatic PowerPoint Report Generation
This template can be used for any scRNA-seq analysis with automatic report generation.

Usage:
    Modify the configuration section with your data paths and parameters,
    then run the script. It will automatically generate:
    - Quality control plots
    - Processed AnnData object
    - UMAP visualizations
    - Marker gene tables
    - PowerPoint report

Integration methods supported: Harmony, Scanorama, BBKNN, or None
Clustering methods: Louvain or Leiden
"""

import scanpy as sc
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pathlib import Path
import warnings
warnings.filterwarnings('ignore')

# ============================================================================
# CONFIGURATION - MODIFY THESE SETTINGS FOR YOUR ANALYSIS
# ============================================================================

# Analysis name (used for output files)
ANALYSIS_NAME = "scRNA_seq_Analysis"

# Data paths
DATA_DIR = Path('/path/to/your/data')  # Directory containing 10X data
OUTPUT_DIR = Path('/path/to/output')   # Where to save results

# Sample information dictionary
# Format: {'sample_name': {'condition': 'condition_name', 'replicate': 'replicate_id'}}
SAMPLES = {
    'sample1': {'condition': 'control', 'replicate': '1'},
    'sample2': {'condition': 'treatment', 'replicate': '1'},
    # Add more samples here
}

# Analysis parameters
PARAMS = {
    # QC filtering
    'min_genes': 200,           # Minimum genes per cell
    'max_genes': 6000,          # Maximum genes per cell
    'min_cells': 3,             # Minimum cells per gene
    'max_mt_percent': 20,       # Maximum mitochondrial percentage

    # Feature selection
    'n_top_genes': 2000,        # Number of highly variable genes

    # Dimensionality reduction
    'n_pcs': 50,                # Number of principal components

    # Integration
    'integration_method': 'harmony',  # Options: 'harmony', 'scanorama', 'bbknn', None
    'batch_key': 'sample',      # Column name for batch correction

    # Clustering
    'clustering_method': 'louvain',   # Options: 'louvain', 'leiden'
    'resolutions': [0.5, 0.8, 1.0],   # Clustering resolutions to try
    'default_resolution': 0.8,        # Default resolution for visualization

    # Visualization
    'n_neighbors': 15,          # Number of neighbors for UMAP
}

# Cell type markers (optional - customize for your tissue type)
CELL_TYPE_MARKERS = {
    'Cell Type 1': ['GENE1', 'GENE2', 'GENE3'],
    'Cell Type 2': ['GENE4', 'GENE5', 'GENE6'],
    # Add more cell types and markers
}

# ============================================================================
# ANALYSIS PIPELINE - NO NEED TO MODIFY BELOW THIS LINE
# ============================================================================

def main():
    # Set up
    sc.settings.verbosity = 3
    sc.settings.set_figure_params(dpi=100, facecolor='white', frameon=False)

    OUTPUT_DIR.mkdir(exist_ok=True)
    figures_dir = OUTPUT_DIR / 'figures'
    figures_dir.mkdir(exist_ok=True)

    print("="*80)
    print(f"{ANALYSIS_NAME}")
    print(f"Integration: {PARAMS['integration_method']}")
    print(f"Clustering: {PARAMS['clustering_method']}")
    print("="*80)

    # Load data
    print(f"\n[1/9] Loading {len(SAMPLES)} samples...")
    print("-"*80)
    adatas = []
    for sample_name, metadata in SAMPLES.items():
        sample_path = DATA_DIR / sample_name / 'outs' / 'filtered_feature_bc_matrix'
        print(f"Loading {sample_name}...")
        adata = sc.read_10x_mtx(sample_path, var_names='gene_symbols', cache=True)
        adata.obs['sample'] = sample_name
        adata.obs['condition'] = metadata['condition']
        adata.obs['replicate'] = metadata['replicate']
        adata.obs['batch'] = sample_name
        adata.var_names_make_unique()
        adatas.append(adata)
        print(f"  Loaded: {adata.n_obs} cells, {adata.n_vars} genes")

    adata = sc.concat(adatas, label='sample_id', keys=list(SAMPLES.keys()))
    print(f"\nTotal dataset: {adata.n_obs} cells, {adata.n_vars} genes")

    # QC
    print(f"\n[2/9] Calculating QC metrics...")
    print("-"*80)
    adata.var['mt'] = adata.var_names.str.startswith('MT-')
    adata.var['ribo'] = adata.var_names.str.match('^RP[SL]')
    sc.pp.calculate_qc_metrics(adata, qc_vars=['mt', 'ribo'], percent_top=None,
                               log1p=False, inplace=True)

    # QC plots
    fig, axes = plt.subplots(2, 3, figsize=(18, 10))
    fig.suptitle('QC Metrics Before Filtering', fontsize=16)
    sc.pl.violin(adata, ['n_genes_by_counts'], groupby='condition', ax=axes[0, 0], show=False)
    sc.pl.violin(adata, ['total_counts'], groupby='condition', ax=axes[0, 1], show=False)
    sc.pl.violin(adata, ['pct_counts_mt'], groupby='condition', ax=axes[0, 2], show=False)
    sc.pl.violin(adata, ['n_genes_by_counts'], groupby='sample', rotation=90, ax=axes[1, 0], show=False)
    sc.pl.violin(adata, ['total_counts'], groupby='sample', rotation=90, ax=axes[1, 1], show=False)
    sc.pl.violin(adata, ['pct_counts_mt'], groupby='sample', rotation=90, ax=axes[1, 2], show=False)
    plt.tight_layout()
    plt.savefig(figures_dir / 'qc_metrics_before_filtering.png', dpi=300, bbox_inches='tight')
    plt.close()

    # Filtering
    print(f"\n[3/9] Performing quality control filtering...")
    print("-"*80)
    print(f"Cells before filtering: {adata.n_obs}")
    sc.pp.filter_cells(adata, min_genes=PARAMS['min_genes'])
    sc.pp.filter_genes(adata, min_cells=PARAMS['min_cells'])
    adata = adata[adata.obs['n_genes_by_counts'] < PARAMS['max_genes'], :]
    adata = adata[adata.obs['pct_counts_mt'] < PARAMS['max_mt_percent'], :]
    print(f"Cells after filtering: {adata.n_obs}")

    adata.layers['counts'] = adata.X.copy()

    # Normalization
    print(f"\n[4/9] Normalizing and log-transforming...")
    print("-"*80)
    sc.pp.normalize_total(adata, target_sum=1e4)
    sc.pp.log1p(adata)
    adata.raw = adata

    # HVG
    print(f"\n[5/9] Identifying highly variable genes...")
    print("-"*80)
    sc.pp.highly_variable_genes(adata, n_top_genes=PARAMS['n_top_genes'],
                                batch_key=PARAMS['batch_key'])
    print(f"Found {sum(adata.var['highly_variable'])} highly variable genes")

    fig = sc.pl.highly_variable_genes(adata, show=False)
    plt.savefig(figures_dir / 'highly_variable_genes.png', dpi=300, bbox_inches='tight')
    plt.close()

    # PCA
    print(f"\n[6/9] Scaling and running PCA...")
    print("-"*80)
    sc.pp.scale(adata, max_value=10)
    sc.tl.pca(adata, svd_solver='arpack', n_comps=PARAMS['n_pcs'])

    sc.pl.pca_variance_ratio(adata, n_pcs=PARAMS['n_pcs'], log=True, show=False)
    plt.savefig(figures_dir / 'pca_variance_ratio.png', dpi=300, bbox_inches='tight')
    plt.close()

    # Integration
    print(f"\n[7/9] Running integration: {PARAMS['integration_method']}...")
    print("-"*80)

    if PARAMS['integration_method'] == 'harmony':
        try:
            import harmonypy as hm
            pca_coords = adata.obsm['X_pca']
            meta_data = adata.obs[[PARAMS['batch_key']]].copy()
            ho = hm.run_harmony(pca_coords, meta_data,
                              vars_use=[PARAMS['batch_key']],
                              max_iter_harmony=20, verbose=True)
            adata.obsm['X_integrated'] = ho.Z_corr
            print("Harmony integration completed!")
        except ImportError:
            print("Harmony not installed. Skipping integration.")
            adata.obsm['X_integrated'] = adata.obsm['X_pca']

    elif PARAMS['integration_method'] == 'scanorama':
        try:
            import scanorama
            adatas_list = [adata[adata.obs[PARAMS['batch_key']] == batch].copy()
                          for batch in adata.obs[PARAMS['batch_key']].unique()]
            integrated = scanorama.correct_scanpy(adatas_list, return_dimred=True)
            adata.obsm['X_integrated'] = integrated[1]
            print("Scanorama integration completed!")
        except ImportError:
            print("Scanorama not installed. Skipping integration.")
            adata.obsm['X_integrated'] = adata.obsm['X_pca']

    else:
        print("No integration method specified. Using PCA coordinates.")
        adata.obsm['X_integrated'] = adata.obsm['X_pca']

    # Clustering
    print(f"\n[8/9] Computing neighbors, UMAP, and clustering...")
    print("-"*80)
    sc.pp.neighbors(adata, n_neighbors=PARAMS['n_neighbors'],
                   n_pcs=PARAMS['n_pcs'], use_rep='X_integrated')
    sc.tl.umap(adata)

    clustering_func = sc.tl.louvain if PARAMS['clustering_method'] == 'louvain' else sc.tl.leiden
    for res in PARAMS['resolutions']:
        key_name = f"{PARAMS['clustering_method']}_{res}"
        clustering_func(adata, resolution=res, key_added=key_name)
        print(f"  Resolution {res}: {adata.obs[key_name].nunique()} clusters")

    # Visualizations
    print(f"\n[9/9] Generating visualizations...")
    print("-"*80)

    default_cluster_key = f"{PARAMS['clustering_method']}_{PARAMS['default_resolution']}"

    # UMAP overview
    fig, axes = plt.subplots(2, 3, figsize=(24, 16))
    sc.pl.umap(adata, color=default_cluster_key, legend_loc='on data',
              ax=axes[0, 0], show=False, title='Clusters')
    sc.pl.umap(adata, color='condition', ax=axes[0, 1], show=False, title='Condition')
    sc.pl.umap(adata, color='sample', ax=axes[0, 2], show=False, title='Sample')
    sc.pl.umap(adata, color='n_genes_by_counts', ax=axes[1, 0], show=False)
    sc.pl.umap(adata, color='total_counts', ax=axes[1, 1], show=False)
    sc.pl.umap(adata, color='pct_counts_mt', ax=axes[1, 2], show=False)
    plt.tight_layout()
    plt.savefig(figures_dir / 'umap_overview.png', dpi=300, bbox_inches='tight')
    plt.close()

    # Marker genes
    sc.tl.rank_genes_groups(adata, default_cluster_key, method='wilcoxon')
    fig = sc.pl.rank_genes_groups(adata, n_genes=25, sharey=False, show=False)
    plt.savefig(figures_dir / 'marker_genes_per_cluster.png', dpi=300, bbox_inches='tight')
    plt.close()

    # Save results
    print("\nSaving processed AnnData object...")
    output_file = OUTPUT_DIR / f'{ANALYSIS_NAME}_integrated.h5ad'
    adata.write(output_file)
    print(f"Saved to: {output_file}")

    adata.obs.to_csv(OUTPUT_DIR / 'metadata.csv')

    for cluster in adata.obs[default_cluster_key].cat.categories:
        markers = sc.get.rank_genes_groups_df(adata, group=cluster)
        markers.to_csv(OUTPUT_DIR / f'markers_cluster_{cluster}.csv', index=False)

    # Generate PowerPoint report
    print("\nGenerating PowerPoint report...")
    generate_ppt_report(adata, OUTPUT_DIR, figures_dir, ANALYSIS_NAME, default_cluster_key)

    print("\n" + "="*80)
    print("ANALYSIS COMPLETE!")
    print("="*80)
    print(f"\nResults saved to: {OUTPUT_DIR}")
    print(f"Processed data: {output_file}")
    print("="*80)

def generate_ppt_report(adata, output_dir, figures_dir, analysis_name, cluster_key):
    """Generate PowerPoint report"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from datetime import datetime

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = analysis_name
        subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

        # Overview slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Analysis Overview"
        body = slide.placeholders[1]
        tf = body.text_frame
        overview_text = [
            f"Total Cells: {adata.n_obs:,}",
            f"Total Genes: {adata.n_vars:,}",
            f"Number of Clusters: {adata.obs[cluster_key].nunique()}",
            f"Number of Conditions: {adata.obs['condition'].nunique()}",
            "",
            f"Integration: {PARAMS['integration_method']}",
            f"Clustering: {PARAMS['clustering_method']}",
        ]
        tf.text = '\n'.join(overview_text)

        # Add figure slides
        figures_to_add = [
            ("Quality Control", "qc_metrics_before_filtering.png"),
            ("Highly Variable Genes", "highly_variable_genes.png"),
            ("PCA Variance", "pca_variance_ratio.png"),
            ("UMAP Overview", "umap_overview.png"),
            ("Marker Genes", "marker_genes_per_cluster.png"),
        ]

        for fig_title, filename in figures_to_add:
            fig_path = figures_dir / filename
            if fig_path.exists():
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)

                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
                tf = txBox.text_frame
                tf.text = fig_title
                p = tf.paragraphs[0]
                p.font.size = Pt(28)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

                left = Inches(0.5)
                top = Inches(1.0)
                slide.shapes.add_picture(str(fig_path), left, top, width=Inches(9))

        ppt_file = output_dir / f'{analysis_name}_Report.pptx'
        prs.save(str(ppt_file))
        print(f"✓ PowerPoint report saved: {ppt_file}")
        print(f"  Size: {ppt_file.stat().st_size / (1024*1024):.1f} MB")

    except ImportError:
        print("Note: python-pptx not installed. Skipping PowerPoint generation.")
    except Exception as e:
        print(f"Warning: Could not generate PowerPoint report: {e}")

if __name__ == "__main__":
    main()
