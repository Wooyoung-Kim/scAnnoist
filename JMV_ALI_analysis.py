#!/usr/bin/env python3
"""
Single-cell RNA-seq analysis of JMV ALI organoid data
Sample type: ALI organoid (non-immune)
Integration method: Harmony
"""

import scanpy as sc
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pathlib import Path
import warnings
warnings.filterwarnings('ignore')

# Set up scanpy
sc.settings.verbosity = 3
sc.settings.set_figure_params(dpi=100, facecolor='white', frameon=False)

# Create output directory
output_dir = Path('/mnt2/kwy/scrna_agent/JMV_ALI_results')
output_dir.mkdir(exist_ok=True)
figures_dir = output_dir / 'figures'
figures_dir.mkdir(exist_ok=True)

print("="*80)
print("JMV ALI Organoid scRNA-seq Analysis")
print("Integration method: Harmony")
print("="*80)

# Define sample information
data_dir = Path('/home/kwy7605/data_61/SARS/Count/JMV_ALI')

samples = {
    'JMV_Control-1': {'condition': 'Control', 'replicate': '1'},
    'JMV_Control-2': {'condition': 'Control', 'replicate': '2'},
    'JMV_Control-3': {'condition': 'Control', 'replicate': '3'},
    'JMV_HCoV_OC43-1': {'condition': 'HCoV_OC43', 'replicate': '1'},
    'JMV_HCoV_OC43-2': {'condition': 'HCoV_OC43', 'replicate': '2'},
    'JMV_HCoV_OC43-3': {'condition': 'HCoV_OC43', 'replicate': '3'},
    'JMV_MERS-1': {'condition': 'MERS', 'replicate': '1'},
    'JMV_MERS-2': {'condition': 'MERS', 'replicate': '2'},
    'JMV_MERS-3': {'condition': 'MERS', 'replicate': '3'},
    'JMV_SARS-CoV-1-1': {'condition': 'SARS-CoV-1', 'replicate': '1'},
    'JMV_SARS-CoV-1-2': {'condition': 'SARS-CoV-1', 'replicate': '2'},
    'JMV_SARS-CoV-1-3': {'condition': 'SARS-CoV-1', 'replicate': '3'},
    'JMV_SARS-CoV-2_L-1': {'condition': 'SARS-CoV-2', 'replicate': '1'},
    'JMV_SARS-CoV-2_L-2': {'condition': 'SARS-CoV-2', 'replicate': '2'},
    'JMV_SARS-CoV-2_L-3': {'condition': 'SARS-CoV-2', 'replicate': '3'},
}

print(f"\n[1/9] Loading {len(samples)} samples...")
print("-"*80)

# Load all samples
adatas = []
for sample_name, metadata in samples.items():
    sample_path = data_dir / sample_name / 'outs' / 'filtered_feature_bc_matrix'
    print(f"Loading {sample_name}...")

    adata = sc.read_10x_mtx(sample_path, var_names='gene_symbols', cache=True)

    # Add metadata
    adata.obs['sample'] = sample_name
    adata.obs['condition'] = metadata['condition']
    adata.obs['replicate'] = metadata['replicate']
    adata.obs['batch'] = sample_name  # For Harmony integration

    # Make variable names unique
    adata.var_names_make_unique()

    adatas.append(adata)
    print(f"  Loaded: {adata.n_obs} cells, {adata.n_vars} genes")

# Concatenate all samples
adata = sc.concat(adatas, label='sample_id', keys=list(samples.keys()))
print(f"\nTotal dataset: {adata.n_obs} cells, {adata.n_vars} genes")

# Calculate QC metrics
print(f"\n[2/9] Calculating QC metrics...")
print("-"*80)

# Mitochondrial genes
adata.var['mt'] = adata.var_names.str.startswith('MT-')
# Ribosomal genes
adata.var['ribo'] = adata.var_names.str.match('^RP[SL]')
# Hemoglobin genes
adata.var['hb'] = adata.var_names.str.contains('^HB[^(P)]')

sc.pp.calculate_qc_metrics(
    adata,
    qc_vars=['mt', 'ribo', 'hb'],
    percent_top=None,
    log1p=False,
    inplace=True
)

# Plot QC metrics before filtering
print("Generating QC violin plots...")
fig, axes = plt.subplots(2, 4, figsize=(20, 10))
fig.suptitle('QC Metrics Before Filtering', fontsize=16, y=1.02)

sc.pl.violin(adata, ['n_genes_by_counts'], groupby='condition', ax=axes[0, 0], show=False)
axes[0, 0].set_title('Genes per cell')

sc.pl.violin(adata, ['total_counts'], groupby='condition', ax=axes[0, 1], show=False)
axes[0, 1].set_title('Total counts per cell')

sc.pl.violin(adata, ['pct_counts_mt'], groupby='condition', ax=axes[0, 2], show=False)
axes[0, 2].set_title('% Mitochondrial')

sc.pl.violin(adata, ['pct_counts_ribo'], groupby='condition', ax=axes[0, 3], show=False)
axes[0, 3].set_title('% Ribosomal')

# By sample
sc.pl.violin(adata, ['n_genes_by_counts'], groupby='sample', rotation=90, ax=axes[1, 0], show=False)
axes[1, 0].set_title('Genes per cell by sample')

sc.pl.violin(adata, ['total_counts'], groupby='sample', rotation=90, ax=axes[1, 1], show=False)
axes[1, 1].set_title('Total counts by sample')

sc.pl.violin(adata, ['pct_counts_mt'], groupby='sample', rotation=90, ax=axes[1, 2], show=False)
axes[1, 2].set_title('% Mitochondrial by sample')

sc.pl.violin(adata, ['pct_counts_ribo'], groupby='sample', rotation=90, ax=axes[1, 3], show=False)
axes[1, 3].set_title('% Ribosomal by sample')

plt.tight_layout()
plt.savefig(figures_dir / 'qc_metrics_before_filtering.png', dpi=300, bbox_inches='tight')
plt.close()

# Print summary statistics
print("\nQC Summary Statistics:")
print(f"  Median genes per cell: {np.median(adata.obs['n_genes_by_counts']):.0f}")
print(f"  Median UMIs per cell: {np.median(adata.obs['total_counts']):.0f}")
print(f"  Median % MT: {np.median(adata.obs['pct_counts_mt']):.2f}%")

# Quality control filtering
print(f"\n[3/9] Performing quality control filtering...")
print("-"*80)

print(f"Cells before filtering: {adata.n_obs}")

# Filter cells
sc.pp.filter_cells(adata, min_genes=200)
print(f"After min_genes filter (>200): {adata.n_obs}")

# Filter genes
sc.pp.filter_genes(adata, min_cells=3)
print(f"After min_cells filter (>3): {adata.n_vars} genes")

# Additional filtering based on QC metrics
adata = adata[adata.obs['n_genes_by_counts'] < 6000, :]
print(f"After max genes filter (<6000): {adata.n_obs}")

adata = adata[adata.obs['pct_counts_mt'] < 20, :]
print(f"After MT filter (<20%): {adata.n_obs}")

print(f"\nFinal dataset: {adata.n_obs} cells, {adata.n_vars} genes")

# Save raw counts
adata.layers['counts'] = adata.X.copy()

# Normalization
print(f"\n[4/9] Normalizing and log-transforming...")
print("-"*80)

sc.pp.normalize_total(adata, target_sum=1e4)
sc.pp.log1p(adata)
adata.raw = adata

# Identify highly variable genes
print(f"\n[5/9] Identifying highly variable genes...")
print("-"*80)

sc.pp.highly_variable_genes(adata, n_top_genes=2000, batch_key='sample')
print(f"Found {sum(adata.var['highly_variable'])} highly variable genes")

# Plot highly variable genes
fig = sc.pl.highly_variable_genes(adata, show=False)
plt.savefig(figures_dir / 'highly_variable_genes.png', dpi=300, bbox_inches='tight')
plt.close()

# Scale and PCA
print(f"\n[6/9] Scaling and running PCA...")
print("-"*80)

sc.pp.scale(adata, max_value=10)
sc.tl.pca(adata, svd_solver='arpack', n_comps=50)

# Plot PCA variance ratio
sc.pl.pca_variance_ratio(adata, n_pcs=50, log=True, show=False)
plt.savefig(figures_dir / 'pca_variance_ratio.png', dpi=300, bbox_inches='tight')
plt.close()

# PCA plots before integration
print("Generating PCA plots before integration...")
fig, axes = plt.subplots(1, 2, figsize=(16, 6))
sc.pl.pca(adata, color='condition', ax=axes[0], show=False, title='PCA by condition (before Harmony)')
sc.pl.pca(adata, color='sample', ax=axes[1], show=False, title='PCA by sample (before Harmony)')
plt.tight_layout()
plt.savefig(figures_dir / 'pca_before_harmony.png', dpi=300, bbox_inches='tight')
plt.close()

# Harmony integration
print(f"\n[7/9] Running Harmony integration...")
print("-"*80)

try:
    import harmonypy as hm

    # Run Harmony
    print("Running Harmony on PCA embeddings...")
    print(f"  Batch key: sample")
    print(f"  Number of PCs: 50")

    # Extract PCA coordinates
    pca_coords = adata.obsm['X_pca']

    # Prepare metadata
    meta_data = adata.obs[['sample']].copy()

    # Run Harmony
    ho = hm.run_harmony(
        pca_coords,
        meta_data,
        vars_use=['sample'],
        max_iter_harmony=20,
        verbose=True
    )

    # Store Harmony coordinates
    # Z_corr is already in the correct shape (n_cells x n_components)
    adata.obsm['X_harmony'] = ho.Z_corr

    print("Harmony integration completed successfully!")

except ImportError:
    print("ERROR: harmonypy not installed. Installing now...")
    import subprocess
    subprocess.run(['pip', 'install', 'harmonypy'], check=True)

    import harmonypy as hm

    # Run Harmony
    pca_coords = adata.obsm['X_pca']
    meta_data = adata.obs[['sample']].copy()

    ho = hm.run_harmony(
        pca_coords,
        meta_data,
        vars_use=['sample'],
        max_iter_harmony=20,
        verbose=True
    )

    adata.obsm['X_harmony'] = ho.Z_corr
    print("Harmony integration completed!")

# Clustering and UMAP
print(f"\n[8/9] Computing neighbors, UMAP, and clustering...")
print("-"*80)

# Use Harmony coordinates for neighbors
sc.pp.neighbors(adata, n_neighbors=15, n_pcs=50, use_rep='X_harmony')

# UMAP
print("Computing UMAP...")
sc.tl.umap(adata)

# Louvain clustering
print("Running Louvain clustering...")
sc.tl.louvain(adata, resolution=0.8, key_added='louvain_0.8')
sc.tl.louvain(adata, resolution=1.0, key_added='louvain_1.0')
sc.tl.louvain(adata, resolution=0.5, key_added='louvain_0.5')

print(f"  Resolution 0.5: {adata.obs['louvain_0.5'].nunique()} clusters")
print(f"  Resolution 0.8: {adata.obs['louvain_0.8'].nunique()} clusters")
print(f"  Resolution 1.0: {adata.obs['louvain_1.0'].nunique()} clusters")

# Generate comprehensive visualizations
print(f"\n[9/9] Generating visualizations...")
print("-"*80)

# UMAP plots - overview
fig, axes = plt.subplots(2, 3, figsize=(24, 16))
sc.pl.umap(adata, color='louvain_0.8', legend_loc='on data',
           ax=axes[0, 0], show=False, title='Clusters (resolution=0.8)')
sc.pl.umap(adata, color='condition', ax=axes[0, 1], show=False, title='Condition')
sc.pl.umap(adata, color='sample', ax=axes[0, 2], show=False, title='Sample', legend_fontsize=8)
sc.pl.umap(adata, color='n_genes_by_counts', ax=axes[1, 0], show=False, title='Genes per cell')
sc.pl.umap(adata, color='total_counts', ax=axes[1, 1], show=False, title='Total counts')
sc.pl.umap(adata, color='pct_counts_mt', ax=axes[1, 2], show=False, title='% Mitochondrial')
plt.tight_layout()
plt.savefig(figures_dir / 'umap_overview.png', dpi=300, bbox_inches='tight')
plt.close()

# UMAP by condition - faceted
print("Creating faceted UMAP by condition...")
conditions = adata.obs['condition'].unique()
n_conditions = len(conditions)
fig, axes = plt.subplots(1, n_conditions, figsize=(6*n_conditions, 5))
if n_conditions == 1:
    axes = [axes]

for idx, condition in enumerate(sorted(conditions)):
    adata_subset = adata[adata.obs['condition'] == condition]
    sc.pl.umap(adata_subset, color='louvain_0.8', legend_loc='on data',
               ax=axes[idx], show=False, title=f'{condition} (n={adata_subset.n_obs})')
plt.tight_layout()
plt.savefig(figures_dir / 'umap_by_condition_faceted.png', dpi=300, bbox_inches='tight')
plt.close()

# Compare PCA before and after Harmony
print("Comparing embeddings before and after Harmony...")
fig, axes = plt.subplots(2, 2, figsize=(16, 16))

# Compute UMAP on original PCA (before Harmony)
adata.obsm['X_pca_umap'] = adata.obsm['X_umap'].copy()
sc.pp.neighbors(adata, n_neighbors=15, n_pcs=50, use_rep='X_pca', key_added='neighbors_pca')
sc.tl.umap(adata, neighbors_key='neighbors_pca')
adata.obsm['X_umap_before_harmony'] = adata.obsm['X_umap'].copy()
adata.obsm['X_umap'] = adata.obsm['X_pca_umap'].copy()

# Plot comparison
sc.pl.embedding(adata, basis='umap_before_harmony', color='condition',
                ax=axes[0, 0], show=False, title='Before Harmony - by condition')
sc.pl.embedding(adata, basis='umap_before_harmony', color='sample',
                ax=axes[0, 1], show=False, title='Before Harmony - by sample', legend_fontsize=8)
sc.pl.umap(adata, color='condition', ax=axes[1, 0], show=False, title='After Harmony - by condition')
sc.pl.umap(adata, color='sample', ax=axes[1, 1], show=False, title='After Harmony - by sample', legend_fontsize=8)
plt.tight_layout()
plt.savefig(figures_dir / 'harmony_comparison.png', dpi=300, bbox_inches='tight')
plt.close()

# Cell type markers for airway epithelial cells
print("Visualizing cell type markers...")

airway_markers = {
    'Basal cells': ['TP63', 'KRT5', 'KRT14'],
    'Ciliated cells': ['FOXJ1', 'PIFO', 'RSPH1'],
    'Club/Secretory cells': ['SCGB1A1', 'SCGB3A2', 'MUC5AC'],
    'Goblet cells': ['MUC5AC', 'MUC5B', 'TFF3'],
    'Ionocytes': ['FOXI1', 'CFTR', 'ATP6V1B1'],
    'Neuroendocrine': ['CHGA', 'CHGB', 'ASCL1'],
}

# Check which markers are present
available_markers = []
for cell_type, markers in airway_markers.items():
    for marker in markers:
        if marker in adata.var_names:
            available_markers.append(marker)

print(f"Found {len(available_markers)} cell type markers in the dataset")

if len(available_markers) > 0:
    # Plot top markers
    markers_to_plot = available_markers[:12] if len(available_markers) > 12 else available_markers

    n_markers = len(markers_to_plot)
    n_cols = 4
    n_rows = (n_markers + n_cols - 1) // n_cols

    fig, axes = plt.subplots(n_rows, n_cols, figsize=(6*n_cols, 5*n_rows))
    axes = axes.flatten() if n_markers > 1 else [axes]

    for idx, marker in enumerate(markers_to_plot):
        sc.pl.umap(adata, color=marker, ax=axes[idx], show=False,
                   title=marker, vmax='p99', cmap='Reds')

    # Hide empty subplots
    for idx in range(n_markers, len(axes)):
        axes[idx].axis('off')

    plt.tight_layout()
    plt.savefig(figures_dir / 'cell_type_markers.png', dpi=300, bbox_inches='tight')
    plt.close()

# Viral receptor expression
print("Visualizing viral receptor expression...")
viral_receptors = ['ACE2', 'TMPRSS2', 'DPP4', 'ANPEP']
available_receptors = [r for r in viral_receptors if r in adata.var_names]

if len(available_receptors) > 0:
    fig, axes = plt.subplots(1, len(available_receptors), figsize=(6*len(available_receptors), 5))
    if len(available_receptors) == 1:
        axes = [axes]

    for idx, receptor in enumerate(available_receptors):
        sc.pl.umap(adata, color=receptor, ax=axes[idx], show=False,
                   title=f'{receptor} expression', vmax='p99', cmap='viridis')

    plt.tight_layout()
    plt.savefig(figures_dir / 'viral_receptors.png', dpi=300, bbox_inches='tight')
    plt.close()

# Cluster composition analysis
print("Analyzing cluster composition...")

# Cells per condition
composition_df = pd.crosstab(adata.obs['louvain_0.8'], adata.obs['condition'], normalize='columns') * 100

fig, ax = plt.subplots(figsize=(12, 8))
composition_df.plot(kind='bar', stacked=True, ax=ax, colormap='tab10')
ax.set_xlabel('Cluster')
ax.set_ylabel('Percentage of cells')
ax.set_title('Cluster composition by condition')
ax.legend(title='Condition', bbox_to_anchor=(1.05, 1), loc='upper left')
plt.xticks(rotation=0)
plt.tight_layout()
plt.savefig(figures_dir / 'cluster_composition_by_condition.png', dpi=300, bbox_inches='tight')
plt.close()

# Cell counts per cluster and condition
counts_df = pd.crosstab(adata.obs['louvain_0.8'], adata.obs['condition'])

fig, ax = plt.subplots(figsize=(12, 8))
counts_df.plot(kind='bar', ax=ax, colormap='tab10')
ax.set_xlabel('Cluster')
ax.set_ylabel('Number of cells')
ax.set_title('Cell counts per cluster and condition')
ax.legend(title='Condition', bbox_to_anchor=(1.05, 1), loc='upper left')
plt.xticks(rotation=0)
plt.tight_layout()
plt.savefig(figures_dir / 'cell_counts_per_cluster.png', dpi=300, bbox_inches='tight')
plt.close()

# Find marker genes
print("Finding marker genes for each cluster...")
sc.tl.rank_genes_groups(adata, 'louvain_0.8', method='wilcoxon', key_added='rank_genes_louvain_0.8')

# Plot top marker genes
fig = sc.pl.rank_genes_groups(adata, n_genes=25, sharey=False, key='rank_genes_louvain_0.8', show=False)
plt.savefig(figures_dir / 'marker_genes_per_cluster.png', dpi=300, bbox_inches='tight')
plt.close()

# Dotplot of top markers
top_markers = []
for cluster in adata.obs['louvain_0.8'].cat.categories:
    markers = sc.get.rank_genes_groups_df(adata, group=cluster, key='rank_genes_louvain_0.8')
    top_markers.extend(markers.head(5)['names'].tolist())

top_markers = list(dict.fromkeys(top_markers))[:50]  # Remove duplicates, keep order, limit to 50

if len(top_markers) > 0:
    fig = sc.pl.dotplot(adata, top_markers, groupby='louvain_0.8',
                        show=False, figsize=(16, 8))
    plt.savefig(figures_dir / 'marker_genes_dotplot.png', dpi=300, bbox_inches='tight')
    plt.close()

# Save marker genes to file
print("Saving marker genes to CSV...")
for cluster in adata.obs['louvain_0.8'].cat.categories:
    markers = sc.get.rank_genes_groups_df(adata, group=cluster, key='rank_genes_louvain_0.8')
    markers.to_csv(output_dir / f'markers_cluster_{cluster}.csv', index=False)

# Summary statistics
print("\n" + "="*80)
print("ANALYSIS SUMMARY")
print("="*80)

print("\nDataset Summary:")
print(f"  Total cells: {adata.n_obs}")
print(f"  Total genes: {adata.n_vars}")
print(f"  Number of samples: {adata.obs['sample'].nunique()}")
print(f"  Number of conditions: {adata.obs['condition'].nunique()}")

print("\nCells per condition:")
print(adata.obs['condition'].value_counts().sort_index())

print("\nCells per sample:")
print(adata.obs['sample'].value_counts().sort_index())

print(f"\nClusters (resolution=0.8): {adata.obs['louvain_0.8'].nunique()}")
print("\nCells per cluster:")
print(adata.obs['louvain_0.8'].value_counts().sort_index())

# Save the processed data
print("\nSaving processed AnnData object...")
output_file = output_dir / 'JMV_ALI_harmony_integrated.h5ad'
adata.write(output_file)
print(f"Saved to: {output_file}")

# Save metadata
print("Saving metadata...")
adata.obs.to_csv(output_dir / 'metadata.csv')

# Generate PowerPoint report
print("\nGenerating PowerPoint report...")
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from datetime import datetime

    def create_title_slide(prs, title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = subtitle
        return slide

    def create_section_header(prs, title):
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        return slide

    def add_image_to_slide(prs, title, image_path, width=None):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        if image_path.exists():
            left = Inches(0.5)
            top = Inches(1.0)
            slide.shapes.add_picture(str(image_path), left, top, width=width if width else Inches(9))
        return slide

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(prs, "JMV ALI Organoid scRNA-seq Analysis",
                      f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    # Overview slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Analysis Overview"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    overview_text = [
        f"Total Cells: {adata.n_obs:,}",
        f"Total Genes: {adata.n_vars:,}",
        f"Number of Clusters: {adata.obs['louvain_0.8'].nunique()}",
        f"Number of Conditions: {adata.obs['condition'].nunique()}",
        "",
        "Pipeline: QC → Normalization → Harmony → Clustering → Marker Detection"
    ]
    tf.text = '\n'.join(overview_text)

    # Add key figures
    figures_to_add = [
        ("Quality Control", "qc_metrics_before_filtering.png"),
        ("Harmony Integration", "harmony_comparison.png"),
        ("UMAP Overview", "umap_overview.png"),
        ("Summary", "summary_overview.png"),
        ("Cell Type Annotation", "umap_cell_type_annotation.png"),  # Annotation UMAP
        ("Cell Type Markers", "cell_type_markers.png"),
        ("Viral Receptors", "viral_receptors.png"),
        ("Marker Genes", "marker_genes_per_cluster.png"),
    ]

    for title, filename in figures_to_add:
        fig_path = figures_dir / filename
        if fig_path.exists():
            add_image_to_slide(prs, title, fig_path)

    # Save PPT
    ppt_file = output_dir / 'JMV_ALI_Analysis_Report.pptx'
    prs.save(str(ppt_file))
    print(f"✓ PowerPoint report saved: {ppt_file}")
    print(f"  Size: {ppt_file.stat().st_size / (1024*1024):.1f} MB")

except ImportError:
    print("Note: python-pptx not installed. Skipping PowerPoint generation.")
    print("      Install with: pip install python-pptx")
except Exception as e:
    print(f"Warning: Could not generate PowerPoint report: {e}")

print("\n" + "="*80)
print("ANALYSIS COMPLETE!")
print("="*80)
print(f"\nResults saved to: {output_dir}")
print(f"Figures saved to: {figures_dir}")
print(f"\nProcessed data: {output_file}")
print("\nGenerated files:")
print(f"  - QC metrics plots")
print(f"  - Highly variable genes plot")
print(f"  - PCA variance ratio plot")
print(f"  - Harmony integration comparison")
print(f"  - UMAP visualizations")
print(f"  - Cell type marker plots")
print(f"  - Viral receptor expression plots")
print(f"  - Cluster composition plots")
print(f"  - Marker gene plots and tables")
print(f"  - Metadata and marker gene CSV files")
print(f"  - PowerPoint report (PPTX)")
print("="*80)
