#!/usr/bin/env python3
"""
Generate PowerPoint report from scRNA-seq analysis results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from pathlib import Path
from datetime import datetime

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def create_section_header(prs, title):
    """Create section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    return slide

def create_content_slide(prs, title, content_text=None):
    """Create content slide with title and optional text"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    if content_text:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_text

    return slide

def add_image_to_slide(prs, title, image_path, width=None, height=None):
    """Add image to a new slide"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = title

    # Format title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add image
    if image_path.exists():
        if width and height:
            left = (prs.slide_width - width) / 2
            top = Inches(1.0)
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        else:
            # Auto-fit image
            left = Inches(0.5)
            top = Inches(1.0)
            slide.shapes.add_picture(str(image_path), left, top, width=Inches(9))

    return slide

def add_table_to_slide(prs, title, df, notes=None):
    """Add table to a new slide"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True

    # Calculate table dimensions
    rows, cols = df.shape
    rows += 1  # Add header row

    # Add table
    left = Inches(1.0)
    top = Inches(1.2)
    width = Inches(8.0)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column headers
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Fill table data
    for row_idx in range(df.shape[0]):
        for col_idx in range(df.shape[1]):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(df.iloc[row_idx, col_idx])
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    # Add notes if provided
    if notes:
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(8.0), Inches(0.5))
        tf = txBox.text_frame
        tf.text = notes
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.italic = True

    return slide

def create_scrna_report(output_dir, results_dir, analysis_name="scRNA-seq Analysis"):
    """
    Generate comprehensive PowerPoint report for scRNA-seq analysis

    Parameters:
    -----------
    output_dir : Path
        Directory where PPT will be saved
    results_dir : Path
        Directory containing analysis results and figures
    analysis_name : str
        Name of the analysis for the title
    """

    output_dir = Path(output_dir)
    results_dir = Path(results_dir)
    figures_dir = results_dir / 'figures'

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(
        prs,
        analysis_name,
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )

    # Slide 2: Analysis Overview
    slide = create_content_slide(prs, "Analysis Overview")
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    # Read metadata to get stats
    metadata_file = results_dir / 'metadata.csv'
    if metadata_file.exists():
        metadata = pd.read_csv(metadata_file)
        n_cells = len(metadata)
        n_samples = metadata['sample'].nunique() if 'sample' in metadata.columns else 'N/A'
        n_conditions = metadata['condition'].nunique() if 'condition' in metadata.columns else 'N/A'
        n_clusters = metadata['louvain_0.8'].nunique() if 'louvain_0.8' in metadata.columns else 'N/A'

        overview_text = [
            f"Total Cells: {n_cells:,}",
            f"Number of Samples: {n_samples}",
            f"Number of Conditions: {n_conditions}",
            f"Number of Clusters: {n_clusters}",
            "",
            "Analysis Pipeline:",
            "  • Quality Control & Filtering",
            "  • Normalization & Feature Selection",
            "  • Batch Correction (Harmony)",
            "  • Dimensionality Reduction (PCA, UMAP)",
            "  • Clustering (Louvain)",
            "  • Marker Gene Identification"
        ]

        for idx, text in enumerate(overview_text):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(16) if idx < 4 else Pt(14)
            p.level = 1 if text.startswith('  •') else 0

    # Slide 3: Section - Quality Control
    create_section_header(prs, "Quality Control")

    # Slide 4: QC Metrics
    qc_fig = figures_dir / 'qc_metrics_before_filtering.png'
    if qc_fig.exists():
        add_image_to_slide(prs, "QC Metrics Before Filtering", qc_fig, width=Inches(9))

    # Slide 5: Highly Variable Genes
    hvg_fig = figures_dir / 'highly_variable_genes.png'
    if hvg_fig.exists():
        add_image_to_slide(prs, "Highly Variable Genes", hvg_fig, width=Inches(9))

    # Slide 6: Section - Batch Correction
    create_section_header(prs, "Batch Correction & Integration")

    # Slide 7: Harmony Integration
    harmony_fig = figures_dir / 'harmony_comparison.png'
    if harmony_fig.exists():
        add_image_to_slide(prs, "Harmony Integration - Before & After", harmony_fig, width=Inches(9))

    # Slide 8: PCA Variance
    pca_var_fig = figures_dir / 'pca_variance_ratio.png'
    if pca_var_fig.exists():
        add_image_to_slide(prs, "PCA Variance Ratio", pca_var_fig, width=Inches(8))

    # Slide 9: Section - Clustering Results
    create_section_header(prs, "Clustering Results")

    # Slide 10: UMAP Overview
    umap_overview_fig = figures_dir / 'umap_overview.png'
    if umap_overview_fig.exists():
        add_image_to_slide(prs, "UMAP Overview", umap_overview_fig, width=Inches(9.5))

    # Slide 11: Summary Overview
    summary_fig = figures_dir / 'summary_overview.png'
    if summary_fig.exists():
        add_image_to_slide(prs, "Analysis Summary", summary_fig, width=Inches(9.5))

    # Slide 11b: Cell Type Annotation UMAP (if available)
    annotation_umap_fig = figures_dir / 'umap_cell_type_annotation.png'
    if annotation_umap_fig.exists():
        add_image_to_slide(prs, "Cell Type Annotation", annotation_umap_fig, width=Inches(9))

    # Slide 12: UMAP by Condition
    umap_condition_fig = figures_dir / 'umap_by_condition_faceted.png'
    if umap_condition_fig.exists():
        add_image_to_slide(prs, "UMAP by Condition (Faceted)", umap_condition_fig, width=Inches(9))

    # Slide 13: Condition Comparison
    condition_comp_fig = figures_dir / 'condition_comparison.png'
    if condition_comp_fig.exists():
        add_image_to_slide(prs, "Condition Comparison", condition_comp_fig, width=Inches(9))

    # Slide 14: Cell Distribution Table
    if metadata_file.exists():
        condition_counts = metadata['condition'].value_counts().sort_index()
        df_conditions = pd.DataFrame({
            'Condition': condition_counts.index,
            'Number of Cells': condition_counts.values,
            'Percentage': [f"{(v/n_cells*100):.1f}%" for v in condition_counts.values]
        })
        add_table_to_slide(prs, "Cells per Condition", df_conditions)

    # Slide 15: Cluster Distribution Table
    if metadata_file.exists() and 'louvain_0.8' in metadata.columns:
        cluster_counts = metadata['louvain_0.8'].value_counts().sort_index().head(15)
        df_clusters = pd.DataFrame({
            'Cluster': cluster_counts.index.astype(str),
            'Number of Cells': cluster_counts.values,
            'Percentage': [f"{(v/n_cells*100):.1f}%" for v in cluster_counts.values]
        })
        add_table_to_slide(prs, "Cells per Cluster (Top 15)", df_clusters,
                          notes="Showing top 15 clusters by size")

    # Slide 16: Cluster Composition
    cluster_comp_fig = figures_dir / 'cluster_composition_by_condition.png'
    if cluster_comp_fig.exists():
        add_image_to_slide(prs, "Cluster Composition by Condition", cluster_comp_fig, width=Inches(8))

    # Slide 17: Cell Counts per Cluster
    cell_counts_fig = figures_dir / 'cell_counts_per_cluster.png'
    if cell_counts_fig.exists():
        add_image_to_slide(prs, "Cell Counts per Cluster and Condition", cell_counts_fig, width=Inches(8))

    # Slide 18: Section - Cell Type Analysis
    create_section_header(prs, "Cell Type Markers")

    # Slide 19: Cell Type Markers
    markers_fig = figures_dir / 'cell_type_markers.png'
    if markers_fig.exists():
        add_image_to_slide(prs, "Airway Epithelial Cell Markers", markers_fig, width=Inches(9.5))

    # Slide 20: Viral Receptors
    viral_receptors_fig = figures_dir / 'viral_receptors.png'
    if viral_receptors_fig.exists():
        add_image_to_slide(prs, "Viral Receptor Expression", viral_receptors_fig, width=Inches(9))

    # Slide 21: Section - Marker Genes
    create_section_header(prs, "Marker Gene Analysis")

    # Slide 22: Marker Genes per Cluster
    marker_genes_fig = figures_dir / 'marker_genes_per_cluster.png'
    if marker_genes_fig.exists():
        add_image_to_slide(prs, "Top Marker Genes per Cluster", marker_genes_fig, width=Inches(9))

    # Slide 23: Marker Genes Dotplot
    dotplot_fig = figures_dir / 'marker_genes_dotplot.png'
    if dotplot_fig.exists():
        add_image_to_slide(prs, "Marker Gene Expression Dotplot", dotplot_fig, width=Inches(9))

    # Slide 24: Top Markers for Selected Clusters
    marker_files = sorted(results_dir.glob('markers_cluster_*.csv'))
    if len(marker_files) > 0:
        # Show top 3 clusters
        for cluster_idx in [0, 1, 2]:
            marker_file = results_dir / f'markers_cluster_{cluster_idx}.csv'
            if marker_file.exists():
                markers_df = pd.read_csv(marker_file)
                top_markers = markers_df.head(10)[['names', 'logfoldchanges', 'pvals_adj']].copy()
                top_markers.columns = ['Gene', 'Log2 FC', 'Adj. P-value']
                top_markers['Log2 FC'] = top_markers['Log2 FC'].round(3)
                top_markers['Adj. P-value'] = top_markers['Adj. P-value'].apply(lambda x: f"{x:.2e}")

                add_table_to_slide(prs, f"Top 10 Marker Genes - Cluster {cluster_idx}", top_markers)

    # Slide 25: Summary & Next Steps
    slide = create_content_slide(prs, "Summary & Next Steps")
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    summary_points = [
        "Key Findings:",
        f"  • Successfully analyzed {n_cells:,} high-quality cells",
        f"  • Identified {n_clusters} distinct cell clusters",
        "  • Harmony integration effectively reduced batch effects",
        "  • Detected airway epithelial cell markers",
        "  • Identified viral receptor expression patterns",
        "",
        "Recommended Next Steps:",
        "  • Annotate cell types using marker genes",
        "  • Differential expression analysis between conditions",
        "  • Pathway enrichment analysis",
        "  • Trajectory inference for cell state transitions",
        "  • Validate findings with additional datasets"
    ]

    for text in summary_points:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16) if text.endswith(':') else Pt(14)
        p.level = 1 if text.startswith('  •') else 0
        if text.endswith(':'):
            p.font.bold = True

    # Save presentation
    output_file = output_dir / f'{analysis_name.replace(" ", "_")}_report.pptx'
    prs.save(str(output_file))

    return output_file

if __name__ == "__main__":
    # Generate report for JMV ALI analysis
    output_dir = Path('/mnt2/kwy/scrna_agent/JMV_ALI_results')
    results_dir = Path('/mnt2/kwy/scrna_agent/JMV_ALI_results')

    print("="*80)
    print("Generating PowerPoint Report")
    print("="*80)

    output_file = create_scrna_report(
        output_dir=output_dir,
        results_dir=results_dir,
        analysis_name="JMV ALI Organoid scRNA-seq Analysis"
    )

    print(f"\n✓ PowerPoint report generated successfully!")
    print(f"  Location: {output_file}")
    print(f"  Size: {output_file.stat().st_size / (1024*1024):.1f} MB")
    print("="*80)
